"""Extracción de texto plano desde documentos para preview en el modal de procesos.

Los documentos en gta/data/procesos/ son guías que se usaron para construir la
fuente de la verdad (los pasos del proceso). El preview solo necesita mostrar
el contenido legible — no se busca preservar formato.

Estrategia por extensión:
  - .pdf            → pdfplumber
  - .docx / .doc    → mammoth (extrae texto plano + indica tablas/imágenes)
  - .xlsx / .xls    → openpyxl (lista hojas con sus celdas)
  - .pptx / .ppt    → python-pptx (texto de cada slide)
  - .txt / .md      → lectura directa
  - resto           → mensaje "tipo no soportado"

Si la librería falla por documento corrupto, devolvemos el error pero no
crasheamos — el frontend muestra el mensaje y ofrece descargar el archivo.
"""
from __future__ import annotations

from pathlib import Path
from typing import Any, Dict


_DATA_ROOT = Path(__file__).resolve().parents[2] / "data" / "procesos"

_MAX_CHARS = 200_000  # ~200 KB de texto extraído. Suficiente para guías típicas.


def _safe_path(rel_path: str) -> Path:
    """Resuelve path relativo dentro de data/procesos sin permitir escape (..)."""
    if not rel_path:
        raise ValueError("path vacío")
    p = (_DATA_ROOT / rel_path).resolve()
    if not str(p).startswith(str(_DATA_ROOT.resolve())):
        raise ValueError("path fuera de data/procesos")
    if not p.exists():
        raise FileNotFoundError(f"no existe: {rel_path}")
    if not p.is_file():
        raise ValueError("no es un archivo")
    return p


def _truncate(text: str) -> Dict[str, Any]:
    if len(text) > _MAX_CHARS:
        return {"text": text[:_MAX_CHARS], "truncated": True, "total_chars": len(text)}
    return {"text": text, "truncated": False, "total_chars": len(text)}


def extraer_texto(rel_path: str) -> Dict[str, Any]:
    """Devuelve {text, truncated, total_chars, kind}. Lanza ValueError si no soportado."""
    path = _safe_path(rel_path)
    ext = path.suffix.lower()

    if ext == ".pdf":
        return {**_extraer_pdf(path), "kind": "pdf"}
    if ext in (".docx", ".doc"):
        return {**_extraer_docx(path), "kind": "word"}
    if ext in (".xlsx", ".xls"):
        return {**_extraer_xlsx(path), "kind": "excel"}
    if ext in (".pptx", ".ppt"):
        return {**_extraer_pptx(path), "kind": "powerpoint"}
    if ext in (".txt", ".md"):
        return {**_truncate(path.read_text(encoding="utf-8", errors="replace")), "kind": "text"}

    raise ValueError(f"extensión no soportada para preview de texto: {ext}")


# ── Extractores específicos ─────────────────────────────────────────────

def _extraer_pdf(path: Path) -> Dict[str, Any]:
    import pdfplumber
    chunks: list[str] = []
    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            txt = page.extract_text() or ""
            if txt.strip():
                chunks.append(f"── Página {i} ──\n{txt}")
            if sum(len(c) for c in chunks) > _MAX_CHARS:
                break
    return _truncate("\n\n".join(chunks).strip())


def _extraer_docx(path: Path) -> Dict[str, Any]:
    # mammoth maneja .docx; .doc legacy lo intenta pero suele fallar.
    if path.suffix.lower() == ".doc":
        return {"text": "(.doc legacy: convertí a .docx para ver el contenido)", "truncated": False, "total_chars": 0}
    import mammoth
    with open(path, "rb") as f:
        result = mammoth.extract_raw_text(f)
    return _truncate(result.value or "")


def _extraer_xlsx(path: Path) -> Dict[str, Any]:
    if path.suffix.lower() == ".xls":
        return {"text": "(.xls legacy: convertí a .xlsx para ver el contenido)", "truncated": False, "total_chars": 0}
    from openpyxl import load_workbook
    wb = load_workbook(filename=str(path), read_only=True, data_only=True)
    chunks: list[str] = []
    for ws in wb.worksheets:
        chunks.append(f"── Hoja: {ws.title} ──")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                chunks.append("\t".join(cells))
            if sum(len(c) for c in chunks) > _MAX_CHARS:
                break
        if sum(len(c) for c in chunks) > _MAX_CHARS:
            break
    return _truncate("\n".join(chunks))


def _extraer_pptx(path: Path) -> Dict[str, Any]:
    if path.suffix.lower() == ".ppt":
        return {"text": "(.ppt legacy: convertí a .pptx para ver el contenido)", "truncated": False, "total_chars": 0}
    from pptx import Presentation
    prs = Presentation(str(path))
    chunks: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        chunks.append(f"── Slide {i} ──")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        chunks.append(text)
        if sum(len(c) for c in chunks) > _MAX_CHARS:
            break
    return _truncate("\n".join(chunks))


# ── Helper de tipo MIME para preview nativo (PDF, imágenes) ─────────────

def detectar_render_mode(rel_path: str) -> Dict[str, Any]:
    """Para el frontend: ¿se puede renderizar nativo o hay que extraer texto?

    Devuelve {mode, mime?}:
      mode='iframe'  → PDF: el browser lo renderiza con <iframe src=download_url>
      mode='image'   → imagen: <img>
      mode='text'    → texto extraído por el endpoint /preview (Word, Excel, etc.)
      mode='download'→ no se puede previsualizar, solo descargar
    """
    ext = Path(rel_path).suffix.lower()
    if ext == ".pdf":
        return {"mode": "iframe", "mime": "application/pdf"}
    if ext in (".png", ".jpg", ".jpeg", ".gif", ".webp"):
        return {"mode": "image", "mime": f"image/{'jpeg' if ext == '.jpg' else ext.lstrip('.')}"}
    if ext in (".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt", ".txt", ".md"):
        return {"mode": "text"}
    return {"mode": "download"}
